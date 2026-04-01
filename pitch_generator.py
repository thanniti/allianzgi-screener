from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── 1. FUND DATA ──────────────────────────────────────────────────────────────
# Swap this dict for a live data pull later (Yahoo Finance, Bloomberg API, etc.)

fund = {
    "name": "AllianzGI Global MA Blend",
    "type": "Growth",
    "region": "Global",
    "return_1y": 12.8,
    "sharpe": 0.55,
    "max_drawdown": -16.7,
    "esg": 68,
    "sortino_est": round(0.55 * 1.3, 2),
    "allocation": {"Equity": 70, "Bond": 20, "Alternatives": 10},
    "score": 63,
    "date": datetime.date.today().strftime("%B %Y"),
}

# ── 2. COLOUR PALETTE ─────────────────────────────────────────────────────────

NAVY       = RGBColor(0x04, 0x2C, 0x53)   # header background
BLUE_MID   = RGBColor(0x0C, 0x44, 0x7C)   # KPI card background
BLUE_LIGHT = RGBColor(0xB5, 0xD4, 0xF4)   # header text
BLUE_TEXT  = RGBColor(0x18, 0x5F, 0xA5)   # accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG    = RGBColor(0xF1, 0xEF, 0xE8)   # section backgrounds
GRAY_TEXT  = RGBColor(0x5F, 0x5E, 0x5A)
GREEN      = RGBColor(0x1D, 0x9E, 0x75)   # positive / bond
AMBER      = RGBColor(0xBA, 0x75, 0x17)   # neutral / warning
RED        = RGBColor(0xE2, 0x4B, 0x4A)   # negative / drawdown
EQUITY_COL = RGBColor(0x37, 0x8A, 0xDD)
ALT_COL    = RGBColor(0xD8, 0x5A, 0x30)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x1A)

# ── 3. HELPERS ────────────────────────────────────────────────────────────────

def add_textbox(slide, text, left, top, width, height,
                font_size=11, bold=False, color=DARK_TEXT,
                align=PP_ALIGN.LEFT, bg=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if bg:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg
    return txBox

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def kpi_color(metric, value):
    if metric == "return_1y":
        return GREEN if value >= 8 else AMBER
    if metric == "sharpe":
        return GREEN if value >= 0.7 else AMBER if value >= 0.5 else RED
    if metric == "max_drawdown":
        return RED if value <= -15 else AMBER if value <= -10 else GREEN
    return BLUE_LIGHT

# ── 4. SLIDE BUILDER ──────────────────────────────────────────────────────────

def build_pitch_slide(fund: dict, output_path: str = "pitch_one_pager.pptx"):

    prs = Presentation()
    prs.slide_width  = Inches(13.33)   # widescreen 16:9
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # ── HEADER BACKGROUND ────────────────────────────────────────────────────
    add_rect(slide, 0, 0, 13.33, 2.6, NAVY)

    # Fund name
    add_textbox(slide, fund["name"],
                0.3, 0.2, 8, 0.5,
                font_size=22, bold=True, color=BLUE_LIGHT)

    # Subtitle
    sub = f"{fund['type']} · {fund['region']} · Multi-Asset"
    add_textbox(slide, sub,
                0.3, 0.65, 8, 0.3,
                font_size=12, color=RGBColor(0x85, 0xB7, 0xEB))

    # Composite score badge (top right)
    add_rect(slide, 10.5, 0.22, 2.5, 0.35, BLUE_TEXT)
    add_textbox(slide, f"Composite score: {fund['score']} / 100",
                10.5, 0.22, 2.5, 0.35,
                font_size=10, bold=True, color=BLUE_LIGHT,
                align=PP_ALIGN.CENTER)

    # ── KPI CARDS ────────────────────────────────────────────────────────────
    kpis = [
        ("1Y return",    f"{fund['return_1y']}%",   "return_1y",    fund["return_1y"]),
        ("Sharpe ratio", f"{fund['sharpe']}",        "sharpe",       fund["sharpe"]),
        ("Max drawdown", f"{fund['max_drawdown']}%", "max_drawdown", fund["max_drawdown"]),
        ("ESG score",    f"{fund['esg']} / 100",     "esg",          fund["esg"]),
    ]

    kpi_left = 0.3
    kpi_width = 3.0
    kpi_gap = 0.1
    kpi_top = 1.1

    for i, (label, val, metric, raw) in enumerate(kpis):
        x = kpi_left + i * (kpi_width + kpi_gap)
        add_rect(slide, x, kpi_top, kpi_width, 1.1, BLUE_MID)
        add_textbox(slide, label.upper(),
                    x + 0.1, kpi_top + 0.08, kpi_width - 0.2, 0.25,
                    font_size=9, color=RGBColor(0x85, 0xB7, 0xEB))
        add_textbox(slide, val,
                    x + 0.1, kpi_top + 0.35, kpi_width - 0.2, 0.6,
                    font_size=24, bold=True,
                    color=kpi_color(metric, raw))

    # ── LEFT COLUMN: HIGHLIGHTS + RATIONALE ──────────────────────────────────
    col_l = 0.3
    col_w = 7.5
    row_y = 2.8

    add_textbox(slide, "KEY HIGHLIGHTS",
                col_l, row_y, col_w, 0.25,
                font_size=9, color=GRAY_TEXT)

    equity_pct = fund['allocation'].get('Equity', 0)
    bond_pct = fund['allocation'].get('Bond', 0)
    alternatives_pct = fund['allocation'].get('Alternatives', fund['allocation'].get('Cash', 0))

    highlights = [
        f"Strong 1Y return of {fund['return_1y']}% reflects equity-led positioning, outperforming a typical 60/40 benchmark.",
        f"Equity overweight ({equity_pct}%) has been the primary return driver, benefiting from global market recovery.",
        f"{alternatives_pct}% alternatives sleeve provides tail-risk mitigation and low-correlation diversification." if alternatives_pct > 0 else "Flexible allocation across non-equity exposure supports drawdown management.",
        f"ESG score of {fund['esg']} reflects integration of sustainability factors, above median for growth-oriented peers.",
    ]

    for j, h in enumerate(highlights):
        y = row_y + 0.3 + j * 0.38
        add_rect(slide, col_l, y + 0.1, 0.07, 0.07, EQUITY_COL)   # bullet dot
        add_textbox(slide, h,
                    col_l + 0.15, y, col_w - 0.2, 0.35,
                    font_size=10, color=DARK_TEXT)

    rat_y = row_y + 0.3 + len(highlights) * 0.38 + 0.15
    add_textbox(slide, "INVESTMENT RATIONALE",
                col_l, rat_y, col_w, 0.25,
                font_size=9, color=GRAY_TEXT)

    rationale = (
        f"The fund targets long-term capital growth via dynamic allocation across global equities, "
        f"investment-grade bonds, and liquid alternatives. Its growth tilt suits investors with a 5+ "
        f"year horizon. The current macro backdrop — resilient global growth and AI-driven earnings "
        f"expansion — favours the fund's {fund['allocation']['Equity']}% equity positioning, while the "
        f"{fund['allocation']['Bond']}% bond sleeve acts as a duration buffer."
    )
    add_textbox(slide, rationale,
                col_l, rat_y + 0.3, col_w, 0.9,
                font_size=10, color=GRAY_TEXT)

    # ── RIGHT COLUMN: RISK + ALLOCATION + ESG ────────────────────────────────
    col_r = 8.2
    col_rw = 4.8

    add_textbox(slide, "RISK PROFILE",
                col_r, row_y, col_rw, 0.25,
                font_size=9, color=GRAY_TEXT)

    risk_items = [
        ("Volatility",    "Moderate-high", AMBER),
        ("Max drawdown",  f"{fund['max_drawdown']}%", RED),
        ("Sharpe ratio",  str(fund["sharpe"]), AMBER),
        ("Sortino est.",  str(fund["sortino_est"]), GREEN),
    ]

    for k, (rl, rv, rc) in enumerate(risk_items):
        rx = col_r + (k % 2) * 2.35
        ry = row_y + 0.3 + (k // 2) * 0.65
        add_rect(slide, rx, ry, 2.2, 0.55, GRAY_BG)
        add_textbox(slide, rl.upper(),
                    rx + 0.1, ry + 0.04, 2.0, 0.2,
                    font_size=8, color=GRAY_TEXT)
        add_textbox(slide, rv,
                    rx + 0.1, ry + 0.24, 2.0, 0.25,
                    font_size=13, bold=True, color=rc)

    # Risk warning box
    warn_y = row_y + 1.65
    add_rect(slide, col_r, warn_y, col_rw, 0.5, RGBColor(0xFA, 0xEE, 0xDA))
    add_textbox(slide, "Risk note",
                col_r + 0.1, warn_y + 0.04, col_rw - 0.2, 0.2,
                font_size=9, bold=True, color=AMBER)
    add_textbox(slide,
                f"Max drawdown of {fund['max_drawdown']}% warrants monitoring. "
                "Clients should be advised of equity concentration risk.",
                col_r + 0.1, warn_y + 0.22, col_rw - 0.2, 0.25,
                font_size=9, color=RGBColor(0x63, 0x38, 0x06))

    # Asset allocation bar
    alloc_y = warn_y + 0.65
    add_textbox(slide, "ASSET ALLOCATION",
                col_r, alloc_y, col_rw, 0.2,
                font_size=9, color=GRAY_TEXT)

    alloc_colors = {
        "Equity": EQUITY_COL,
        "Bond": GREEN,
        "Alternatives": ALT_COL,
        "Cash": AMBER,
    }
    bar_y = alloc_y + 0.25
    bar_w = col_rw
    x_cursor = col_r
    for asset, pct in fund["allocation"].items():
        seg_w = bar_w * pct / 100
        color = alloc_colors.get(asset, GRAY_BG)
        add_rect(slide, x_cursor, bar_y, seg_w, 0.18, color)
        x_cursor += seg_w

    legend_y = bar_y + 0.25
    for m, (asset, pct) in enumerate(fund["allocation"].items()):
        lx = col_r + m * 1.55
        add_rect(slide, lx, legend_y + 0.04, 0.12, 0.12, alloc_colors[asset])
        add_textbox(slide, f"{asset}  {pct}%",
                    lx + 0.18, legend_y, 1.3, 0.22,
                    font_size=10, color=GRAY_TEXT)

    # ESG bar
    esg_y = legend_y + 0.38
    add_textbox(slide, "ESG INTEGRATION",
                col_r, esg_y, col_rw, 0.2,
                font_size=9, color=GRAY_TEXT)
    add_rect(slide, col_r, esg_y + 0.25, col_rw, 0.14, GRAY_BG)
    add_rect(slide, col_r, esg_y + 0.25, col_rw * fund["esg"] / 100, 0.14, GREEN)
    add_textbox(slide, f"Score: {fund['esg']} / 100",
                col_r, esg_y + 0.42, col_rw, 0.2,
                font_size=9, color=GRAY_TEXT)

    # ── FOOTER ───────────────────────────────────────────────────────────────
    add_rect(slide, 0, 7.1, 13.33, 0.02, GRAY_BG)
    add_textbox(slide,
                "Allianz Global Investors Asia Pacific Ltd. · Hong Kong  |  "
                "For professional and institutional investors only. Not for retail distribution.  |  "
                f"As at {fund['date']}  |  Past performance is not indicative of future results.",
                0.3, 7.15, 12.7, 0.3,
                font_size=8, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    prs.save(output_path)
    print(f"Saved: {output_path}", flush=True)

# ── 5. RUN ────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Starting script...", flush=True)
    build_pitch_slide(fund, "AllianzGI_GlobalMABlend_Pitch.pptx")